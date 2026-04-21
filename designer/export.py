"""Designer — export pipeline (PDF, PPTX, HTML, PNG)."""

from __future__ import annotations

import io
import logging
import os
import pathlib
import re
import zipfile
from typing import Any, Optional

from designer.preview import render_page_html
from designer.state import DesignerProject

logger = logging.getLogger(__name__)

_WORKSPACE = pathlib.Path(
    os.environ.get("THOTH_WORKSPACE", pathlib.Path.home() / "Documents" / "Thoth")
)


class ExportedBytes(bytes):
    """Bytes payload annotated with the actual path written to disk."""

    saved_path: pathlib.Path

    def __new__(cls, payload: bytes, saved_path: pathlib.Path):
        obj = super().__new__(cls, payload)
        obj.saved_path = pathlib.Path(saved_path)
        return obj

_CSS_NUMBER_RE = re.compile(r"-?\d+(?:\.\d+)?")
_RENDERED_DOM_EXPORT_SCRIPT = r"""
() => {
    const normalizeText = (value) => (value || '').replace(/\u00a0/g, ' ').replace(/\s+/g, ' ').trim();
    const parseNumber = (value) => {
        const parsed = Number.parseFloat(value || '0');
        return Number.isFinite(parsed) ? parsed : 0;
    };
    const parseAlpha = (value) => {
        if (!value || value === 'transparent') {
            return 0;
        }
        const text = String(value).trim().toLowerCase();
        if (text.startsWith('#')) {
            return 1;
        }
        const parts = text.match(/[\d.]+/g) || [];
        if (parts.length < 3) {
            return 0;
        }
        return parts.length >= 4 ? Math.max(0, Math.min(1, Number.parseFloat(parts[3]) || 0)) : 1;
    };
    const transformText = (value, transform) => {
        if (!value) {
            return '';
        }
        switch ((transform || '').toLowerCase()) {
            case 'uppercase':
                return value.toUpperCase();
            case 'lowercase':
                return value.toLowerCase();
            case 'capitalize':
                return value.replace(/\b(\w)/g, (match) => match.toUpperCase());
            default:
                return value;
        }
    };
        const INLINE_TAGS = new Set(['span', 'em', 'strong', 'b', 'i', 'u', 'small', 'mark', 'sup', 'sub', 'code']);
    const ownText = (element) => Array.from(element.childNodes)
        .filter((node) => node.nodeType === Node.TEXT_NODE)
        .map((node) => node.textContent || '')
        .join(' ');
        const isInlineLike = (tag, style) => {
            const display = (style?.display || '').toLowerCase();
            return INLINE_TAGS.has(tag) || display === 'inline' || display === 'contents';
        };
    const normalizedVisibleText = (element, style = null) => {
        const resolvedStyle = style || getComputedStyle(element);
        return normalizeText(transformText(element.innerText || '', resolvedStyle.textTransform));
    };
    const ownVisibleText = (element, style = null) => {
        const resolvedStyle = style || getComputedStyle(element);
        return normalizeText(transformText(ownText(element), resolvedStyle.textTransform));
    };
    const textBearingChildElements = (element) => Array.from(element.children).filter((child) => {
        if (!(child instanceof HTMLElement)) {
            return false;
        }
        return !!normalizedVisibleText(child);
    });
    const hasBlockTextDescendant = (element) => Array.from(element.querySelectorAll('*')).some((child) => {
        if (!(child instanceof HTMLElement)) {
            return false;
        }
        const childText = normalizedVisibleText(child);
        if (!childText) {
            return false;
        }
        const childStyle = getComputedStyle(child);
        return !isInlineLike(child.tagName.toLowerCase(), childStyle);
    });
    const ownsLeafTextCluster = (element, style = null) => {
        const resolvedStyle = style || getComputedStyle(element);
        if (hasBlockTextDescendant(element)) {
            return false;
        }
        return !!ownVisibleText(element, resolvedStyle) || textBearingChildElements(element).length <= 1;
    };
        const exportedText = (element, tag, style) => {
        const visibleText = normalizedVisibleText(element, style);
            if (!visibleText) {
                return '';
            }

        const ownTextValue = ownVisibleText(element, style);
        const blockTextDescendant = hasBlockTextDescendant(element);

            if (!isInlineLike(tag, style)) {
            if (!blockTextDescendant) {
                return ownsLeafTextCluster(element, style) ? visibleText : '';
            }
            return ownTextValue;
            }

            const parent = element.parentElement;
            if (!parent) {
                return visibleText;
            }

            const parentStyle = getComputedStyle(parent);
        const parentText = normalizedVisibleText(parent, parentStyle);
        return parentText && ownsLeafTextCluster(parent, parentStyle) ? '' : visibleText;
        };
    const isVisible = (style, rect) => {
        if (!style || style.display === 'none' || style.visibility === 'hidden') {
            return false;
        }
        if (parseFloat(style.opacity || '1') <= 0.01) {
            return false;
        }
        if (rect.width < 4 || rect.height < 4) {
            return false;
        }
        if (rect.right <= 0 || rect.bottom <= 0) {
            return false;
        }
        return true;
    };

    const items = [];
    let order = 0;
    let screenshotCounter = 0;

    for (const element of Array.from(document.body.querySelectorAll('*'))) {
        const rect = element.getBoundingClientRect();
        const style = getComputedStyle(element);
        if (!isVisible(style, rect)) {
            continue;
        }

        const tag = element.tagName.toLowerCase();
    const visibleText = normalizeText(transformText(element.innerText || '', style.textTransform));
    const text = exportedText(element, tag, style);
        const backgroundColor = style.backgroundColor || '';
        const backgroundImage = style.backgroundImage || 'none';
        const borderWidth = Math.max(
            parseNumber(style.borderTopWidth),
            parseNumber(style.borderRightWidth),
            parseNumber(style.borderBottomWidth),
            parseNumber(style.borderLeftWidth),
        );
        const borderColor = style.borderTopColor || style.borderColor || '';
        const borderStyle = style.borderTopStyle || style.borderStyle || 'none';
        const hasFill = parseAlpha(backgroundColor) > 0.02;
        const hasBorder = borderStyle !== 'none' && borderWidth > 0.1 && parseAlpha(borderColor) > 0.02;
        const hasBackgroundImage = backgroundImage !== 'none';
        const zIndexRaw = Number.parseInt(style.zIndex || '0', 10);
        const zIndex = Number.isFinite(zIndexRaw) ? zIndexRaw : 0;

        const base = {
            order: order++,
            tag,
            x: rect.left,
            y: rect.top,
            width: rect.width,
            height: rect.height,
            zIndex,
            backgroundColor,
            backgroundImage,
            borderColor,
            borderWidth,
            borderRadius: style.borderRadius || '0px',
            color: style.color || '',
            fontFamily: style.fontFamily || '',
            fontSize: style.fontSize || '',
            fontStyle: style.fontStyle || 'normal',
            fontWeight: style.fontWeight || '400',
            lineHeight: style.lineHeight || '',
            textAlign: style.textAlign || 'left',
            opacity: parseFloat(style.opacity || '1') || 1,
            paddingTop: style.paddingTop || '0px',
            paddingRight: style.paddingRight || '0px',
            paddingBottom: style.paddingBottom || '0px',
            paddingLeft: style.paddingLeft || '0px',
            display: style.display || 'block',
            alignItems: style.alignItems || 'stretch',
            justifyContent: style.justifyContent || 'flex-start',
            whiteSpace: style.whiteSpace || 'normal',
        };

        if (tag === 'img') {
            const screenshotId = `pptx-node-${screenshotCounter++}`;
            element.setAttribute('data-thoth-pptx-export-id', screenshotId);
            items.push({
                ...base,
                kind: 'image',
                screenshotId,
            });
            continue;
        }

        if (hasBackgroundImage && !visibleText) {
            const screenshotId = `pptx-node-${screenshotCounter++}`;
            element.setAttribute('data-thoth-pptx-export-id', screenshotId);
            items.push({
                ...base,
                kind: 'snapshot',
                screenshotId,
            });
        } else if ((hasFill || hasBorder) && rect.width >= 6 && rect.height >= 6) {
            items.push({
                ...base,
                kind: 'shape',
            });
        }

        if (text) {
            items.push({
                ...base,
                kind: 'text',
                text: tag === 'li' && !text.startsWith('•') ? `• ${text}` : text,
            });
        }
    }

    const bodyStyle = getComputedStyle(document.body);
    const rootStyle = getComputedStyle(document.documentElement);
    const bodyBackground = parseAlpha(bodyStyle.backgroundColor) > 0.02
        ? bodyStyle.backgroundColor
        : rootStyle.backgroundColor;

    return {
        backgroundColor: bodyBackground || '',
        items,
    };
}
"""


def get_export_workspace() -> pathlib.Path:
    """Return the default export workspace and ensure it exists."""
    _WORKSPACE.mkdir(parents=True, exist_ok=True)
    return _WORKSPACE


def _sanitize_name(text: str, max_len: int = 60) -> str:
    safe = "".join(c if c.isalnum() or c in " -_" else "_" for c in (text or ""))
    safe = safe.strip()[:max_len]
    return safe or "Designer Export"


def _resolve_directory(directory: pathlib.Path | str | os.PathLike | None) -> pathlib.Path:
    if directory is None:
        return get_export_workspace()
    resolved = pathlib.Path(directory)
    resolved.mkdir(parents=True, exist_ok=True)
    return resolved


def _parse_page_range(pages_str: Optional[str], total: int) -> list[int]:
    """Parse a page range string like '1-3' or '1,3,5' into 0-based indices."""
    if not pages_str or pages_str.lower() == "all":
        return list(range(total))
    indices = set()
    for part in pages_str.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            a, b = part.split("-", 1)
            start = max(0, int(a.strip()) - 1)
            end = min(total, int(b.strip()))
            indices.update(range(start, end))
        else:
            idx = int(part.strip()) - 1
            if 0 <= idx < total:
                indices.add(idx)
    return sorted(indices) if indices else list(range(total))


def _selected_pages(project: DesignerProject, pages: Optional[str]) -> list[tuple[int, object]]:
    indices = _parse_page_range(pages, len(project.pages))
    return [(i, project.pages[i]) for i in indices if i < len(project.pages)]


def describe_export_destination(
    project: DesignerProject,
    format: str,
    pages: Optional[str] = None,
    mode: Optional[str] = None,
    directory: pathlib.Path | str | os.PathLike | None = None,
) -> pathlib.Path:
    """Return the exact path an export operation will write to."""
    fmt = (format or "").lower().strip()
    out_dir = _resolve_directory(directory)
    safe_name = _sanitize_name(project.name)
    page_count = len(_parse_page_range(pages, len(project.pages)))
    if fmt == "html":
        filename = f"{safe_name}.html"
    elif fmt == "pdf":
        filename = f"{safe_name}.pdf"
    elif fmt == "png":
        filename = f"{safe_name}.png" if page_count <= 1 else f"{safe_name}_pages.zip"
    elif fmt == "pptx":
        filename = f"{safe_name}_editable.pptx" if (mode or "").lower() == "structured" else f"{safe_name}.pptx"
    else:
        raise ValueError(f"Unsupported export format: {format}")
    return out_dir / filename


def _next_available_export_path(path: pathlib.Path) -> pathlib.Path:
    candidate = path
    counter = 2
    while candidate.exists():
        candidate = path.with_name(f"{path.stem} ({counter}){path.suffix}")
        counter += 1
    return candidate


def _permission_denied_message(path: pathlib.Path, label: str) -> str:
    return (
        f"{label} could not be saved to {path} because that file is open or locked. "
        "Close it in PowerPoint or Explorer preview and try again."
    )


def _save_bytes(path: pathlib.Path, data: bytes, label: str) -> bytes:
    path.parent.mkdir(parents=True, exist_ok=True)
    saved_path = path
    try:
        path.write_bytes(data)
    except PermissionError as exc:
        if path.exists():
            fallback_path = _next_available_export_path(path)
            try:
                fallback_path.write_bytes(data)
            except PermissionError as fallback_exc:
                raise PermissionError(_permission_denied_message(path, label)) from fallback_exc
            saved_path = fallback_path
            logger.warning("%s target locked; saved to %s instead of %s", label, saved_path, path)
        else:
            raise PermissionError(_permission_denied_message(path, label)) from exc

    logger.info("Exported %s to %s (%d bytes)", label, saved_path, len(data))
    return ExportedBytes(data, saved_path)


def _escape_html(text: str) -> str:
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def _escape_attr(text: str) -> str:
    return (
        text.replace("&", "&amp;")
        .replace('"', "&quot;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
    )


def build_html_export(project: DesignerProject, pages: Optional[str] = None) -> bytes:
    """Bundle all or selected pages into a single self-contained HTML file."""
    selected = _selected_pages(project, pages)

    from designer.fonts import get_font_css_embedded

    brand = project.brand
    font_families = list(dict.fromkeys([brand.heading_font, brand.body_font] if brand else []))
    embedded_font_css = "\n".join(get_font_css_embedded(f) for f in font_families)

    sections = []
    for i, page in selected:
        html = render_page_html(project, page.html, page_index=i)
        sections.append(
            f'<section id="page-{i}" style="margin-bottom:40px; page-break-after:always;">\n'
            f'<h2 style="font-family:sans-serif;font-size:14px;color:#888;margin-bottom:8px;">'
            f'Page {i + 1}: {page.title}</h2>\n'
            f'<div style="border:1px solid #333;border-radius:8px;overflow:hidden;">\n'
            f'<iframe srcdoc="{_escape_attr(html)}" '
            f'style="width:{project.canvas_width}px;height:{project.canvas_height}px;border:none;" '
            f'sandbox="allow-same-origin"></iframe>\n'
            f'</div>\n</section>'
        )

    indices = _parse_page_range(pages, len(project.pages))
    nav_links = " ".join(
        f'<a href="#page-{i}" style="color:#2563EB;margin-right:12px;">'
        f'{project.pages[i].title}</a>'
        for i in indices if i < len(project.pages)
    )

    full = (
        f"<!DOCTYPE html><html><head><meta charset='utf-8'>"
        f"<title>{_escape_html(project.name)}</title>"
        f"<style>{embedded_font_css}\n"
        f"body{{background:#111;color:#fff;font-family:sans-serif;padding:20px;}}"
        f"a{{text-decoration:none;}}</style></head><body>"
        f"<h1>{_escape_html(project.name)}</h1>"
        f"<nav style='margin-bottom:20px;'>{nav_links}</nav>"
        + "\n".join(sections)
        + "</body></html>"
    )
    return full.encode("utf-8")


def export_html(
    project: DesignerProject,
    pages: Optional[str] = None,
    directory: pathlib.Path | str | os.PathLike | None = None,
) -> bytes:
    data = build_html_export(project, pages)
    out_path = describe_export_destination(project, "html", pages, directory=directory)
    return _save_bytes(out_path, data, "HTML")


def export_pdf(
    project: DesignerProject,
    pages: Optional[str] = None,
    directory: pathlib.Path | str | os.PathLike | None = None,
) -> bytes:
    """Render each page as a PDF using Playwright and merge the result."""
    selected = _selected_pages(project, pages)

    try:
        from playwright.sync_api import sync_playwright
    except ImportError as exc:
        raise RuntimeError(
            "Playwright is required for PDF export. Run: pip install playwright && playwright install chromium"
        ) from exc

    pdf_pages = []
    with sync_playwright() as pw:
        browser = pw.chromium.launch(headless=True)
        for i, page in selected:
            html = render_page_html(project, page.html, page_index=i)
            ctx = browser.new_context(
                viewport={"width": project.canvas_width, "height": project.canvas_height}
            )
            pg = ctx.new_page()
            pg.set_content(html, wait_until="networkidle")
            pdf_pages.append(
                pg.pdf(
                    width=f"{project.canvas_width}px",
                    height=f"{project.canvas_height}px",
                    print_background=True,
                )
            )
            ctx.close()
        browser.close()

    if len(pdf_pages) == 1:
        merged = pdf_pages[0]
    else:
        try:
            from pypdf import PdfReader, PdfWriter
        except ImportError:
            logger.warning("pypdf not installed — returning first page only for multi-page PDF export")
            merged = pdf_pages[0]
        else:
            writer = PdfWriter()
            for pdf_data in pdf_pages:
                reader = PdfReader(io.BytesIO(pdf_data))
                for page in reader.pages:
                    writer.add_page(page)
            buf = io.BytesIO()
            writer.write(buf)
            merged = buf.getvalue()

    out_path = describe_export_destination(project, "pdf", pages, directory=directory)
    return _save_bytes(out_path, merged, "PDF")


def _render_png_screenshots(project: DesignerProject, pages: Optional[str] = None) -> list[tuple[str, bytes]]:
    selected = _selected_pages(project, pages)

    try:
        from playwright.sync_api import sync_playwright
    except ImportError as exc:
        raise RuntimeError("Playwright is required for PNG export.") from exc

    screenshots: list[tuple[str, bytes]] = []
    with sync_playwright() as pw:
        browser = pw.chromium.launch(headless=True)
        for i, page in selected:
            html = render_page_html(project, page.html, page_index=i)
            ctx = browser.new_context(
                viewport={"width": project.canvas_width, "height": project.canvas_height},
            )
            pg = ctx.new_page()
            pg.set_content(html, wait_until="networkidle")
            png_bytes = pg.screenshot(full_page=False, type="png")
            safe_title = _sanitize_name(page.title, max_len=40)
            screenshots.append((f"page_{i + 1}_{safe_title}.png", png_bytes))
            ctx.close()
        browser.close()
    return screenshots


def export_png_files(
    project: DesignerProject,
    pages: Optional[str] = None,
    directory: pathlib.Path | str | os.PathLike | None = None,
) -> list[pathlib.Path]:
    """Write selected pages as individual PNG files and return the paths."""
    screenshots = _render_png_screenshots(project, pages)
    out_dir = _resolve_directory(directory)
    paths: list[pathlib.Path] = []
    for filename, png_bytes in screenshots:
        path = out_dir / filename
        path.write_bytes(png_bytes)
        paths.append(path)
    logger.info("Exported %d individual PNG files to %s", len(paths), out_dir)
    return paths


def export_png(
    project: DesignerProject,
    pages: Optional[str] = None,
    directory: pathlib.Path | str | os.PathLike | None = None,
) -> bytes:
    """Screenshot selected pages as PNG. Returns one PNG or a ZIP for multiple pages."""
    screenshots = _render_png_screenshots(project, pages)
    out_path = describe_export_destination(project, "png", pages, directory=directory)

    if len(screenshots) == 1:
        return _save_bytes(out_path, screenshots[0][1], "PNG")

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as archive:
        for filename, png_bytes in screenshots:
            archive.writestr(filename, png_bytes)
    return _save_bytes(out_path, buf.getvalue(), "PNG ZIP")


def export_pptx_screenshot(
    project: DesignerProject,
    pages: Optional[str] = None,
    directory: pathlib.Path | str | os.PathLike | None = None,
) -> bytes:
    """Render each page as an image via Playwright and embed in PPTX slides."""
    selected = _selected_pages(project, pages)

    try:
        from playwright.sync_api import sync_playwright
    except ImportError as exc:
        raise RuntimeError("Playwright is required for PPTX screenshot export.") from exc
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError as exc:
        raise RuntimeError(
            "python-pptx is required for PPTX export. Run: pip install python-pptx"
        ) from exc

    prs = Presentation()
    prs.slide_width = Emu(int(project.canvas_width * 914400 / 96))
    prs.slide_height = Emu(int(project.canvas_height * 914400 / 96))
    blank_layout = prs.slide_layouts[6]

    with sync_playwright() as pw:
        browser = pw.chromium.launch(headless=True)
        for i, page in selected:
            html = render_page_html(project, page.html, page_index=i)
            ctx = browser.new_context(
                viewport={"width": project.canvas_width, "height": project.canvas_height},
            )
            pg = ctx.new_page()
            pg.set_content(html, wait_until="networkidle")
            png_bytes = pg.screenshot(full_page=False, type="png")
            ctx.close()

            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(io.BytesIO(png_bytes), Emu(0), Emu(0), prs.slide_width, prs.slide_height)
            if page.notes:
                slide.notes_slide.notes_text_frame.text = page.notes
        browser.close()

    buf = io.BytesIO()
    prs.save(buf)
    out_path = describe_export_destination(project, "pptx", pages, directory=directory)
    return _save_bytes(out_path, buf.getvalue(), "PPTX")


def _css_length_to_px(value: Any, default: float = 0.0) -> float:
    if isinstance(value, (int, float)):
        return float(value)
    match = _CSS_NUMBER_RE.search(str(value or ""))
    return float(match.group(0)) if match else default


def _px_to_emu(value: Any) -> int:
    return int(round(max(0.0, _css_length_to_px(value)) * 914400 / 96))


def _px_to_pt(value: Any, default: float = 12.0) -> float:
    pixels = _css_length_to_px(value, default / 0.75)
    return max(1.0, pixels * 72 / 96)


def _parse_css_color(value: Any):
    text = str(value or "").strip().lower()
    if not text or text == "transparent":
        return None
    if text.startswith("#"):
        hex_value = text[1:]
        if len(hex_value) == 3:
            hex_value = "".join(char * 2 for char in hex_value)
        if len(hex_value) == 6:
            return tuple(int(hex_value[i:i + 2], 16) for i in range(0, 6, 2)), 1.0
        return None
    channels = [float(number) for number in re.findall(r"[\d.]+", text)]
    if len(channels) < 3:
        return None
    rgb = tuple(max(0, min(255, int(round(channel)))) for channel in channels[:3])
    alpha = max(0.0, min(1.0, channels[3])) if len(channels) >= 4 else 1.0
    return rgb, alpha


def _font_name_from_css(value: Any) -> str | None:
    families = [family.strip().strip('"\'') for family in str(value or "").split(",")]
    return next((family for family in families if family), None)


def _weight_from_css(value: Any) -> int:
    text = str(value or "").strip().lower()
    if text == "bold":
        return 700
    if text == "normal":
        return 400
    match = _CSS_NUMBER_RE.search(text)
    return int(float(match.group(0))) if match else 400


def _alignment_from_css(value: Any):
    from pptx.enum.text import PP_ALIGN

    normalized = str(value or "left").strip().lower()
    if normalized == "center":
        return PP_ALIGN.CENTER
    if normalized == "right":
        return PP_ALIGN.RIGHT
    if normalized == "justify":
        return PP_ALIGN.JUSTIFY
    return PP_ALIGN.LEFT


def _apply_fill(fill, color_value: Any) -> bool:
    parsed = _parse_css_color(color_value)
    if not parsed:
        return False
    from pptx.dml.color import RGBColor

    rgb, alpha = parsed
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)
    try:
        fill.transparency = max(0.0, min(1.0, 1.0 - alpha))
    except Exception:
        pass
    return True


def _apply_line(shape, item: dict[str, Any]) -> None:
    parsed = _parse_css_color(item.get("borderColor"))
    border_width = _css_length_to_px(item.get("borderWidth"))
    if not parsed or border_width <= 0:
        shape.line.fill.background()
        return
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    rgb, _alpha = parsed
    shape.line.color.rgb = RGBColor(*rgb)
    shape.line.width = Pt(_px_to_pt(border_width))


def _shape_type_for_radius(item: dict[str, Any]):
    from pptx.enum.shapes import MSO_SHAPE

    radius = _css_length_to_px(item.get("borderRadius"))
    return MSO_SHAPE.ROUNDED_RECTANGLE if radius > 0.5 else MSO_SHAPE.RECTANGLE


def _text_vertical_anchor(item: dict[str, Any]):
    from pptx.enum.text import MSO_VERTICAL_ANCHOR

    display = str(item.get("display") or "").strip().lower()
    align_items = str(item.get("alignItems") or "").strip().lower()
    if display == "flex" and align_items in {"center", "flex-end"}:
        return MSO_VERTICAL_ANCHOR.MIDDLE if align_items == "center" else MSO_VERTICAL_ANCHOR.BOTTOM
    return MSO_VERTICAL_ANCHOR.TOP


def _collect_rendered_slide_snapshot(project: DesignerProject, page_html: str, *, page_index: int):
    try:
        from playwright.sync_api import sync_playwright
    except ImportError as exc:
        raise RuntimeError(
            "Playwright is required for editable PPTX export. Run: pip install playwright && playwright install chromium"
        ) from exc

    html = render_page_html(project, page_html, page_index=page_index)
    screenshot_bytes: dict[str, bytes] = {}

    with sync_playwright() as pw:
        browser = pw.chromium.launch(headless=True)
        ctx = browser.new_context(
            viewport={"width": project.canvas_width, "height": project.canvas_height},
            device_scale_factor=1,
        )
        pg = ctx.new_page()
        pg.set_content(html, wait_until="networkidle")
        snapshot = pg.evaluate(_RENDERED_DOM_EXPORT_SCRIPT)

        for item in snapshot.get("items", []):
            screenshot_id = item.get("screenshotId")
            if not screenshot_id:
                continue
            locator = pg.locator(f'[data-thoth-pptx-export-id="{screenshot_id}"]').first
            screenshot_bytes[screenshot_id] = locator.screenshot(animations="disabled", omit_background=True)

        ctx.close()
        browser.close()

    return snapshot, screenshot_bytes


def _rendered_item_sort_key(item: dict[str, Any]) -> tuple[int, int, int]:
    kind_priority = 1 if item.get("kind") == "text" else 0
    return int(item.get("zIndex") or 0), int(item.get("order") or 0), kind_priority


def _add_rendered_item_to_slide(slide, item: dict[str, Any], screenshots: dict[str, bytes]) -> None:
    from pptx.util import Emu, Pt
    from pptx.enum.text import MSO_AUTO_SIZE

    left = Emu(_px_to_emu(item.get("x")))
    top = Emu(_px_to_emu(item.get("y")))
    width = Emu(_px_to_emu(item.get("width")))
    height = Emu(_px_to_emu(item.get("height")))
    kind = item.get("kind")

    if kind in {"image", "snapshot"}:
        screenshot_id = item.get("screenshotId")
        png_bytes = screenshots.get(screenshot_id or "")
        if png_bytes:
            slide.shapes.add_picture(io.BytesIO(png_bytes), left, top, width, height)
        return

    if kind == "shape":
        shape = slide.shapes.add_shape(_shape_type_for_radius(item), left, top, width, height)
        if not _apply_fill(shape.fill, item.get("backgroundColor")):
            shape.fill.background()
        _apply_line(shape, item)
        return

    if kind != "text":
        return

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.word_wrap = str(item.get("whiteSpace") or "normal").strip().lower() not in {"nowrap", "pre"}
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.margin_left = Emu(_px_to_emu(item.get("paddingLeft")))
    text_frame.margin_right = Emu(_px_to_emu(item.get("paddingRight")))
    text_frame.margin_top = Emu(_px_to_emu(item.get("paddingTop")))
    text_frame.margin_bottom = Emu(_px_to_emu(item.get("paddingBottom")))
    text_frame.vertical_anchor = _text_vertical_anchor(item)

    paragraph = text_frame.paragraphs[0]
    paragraph.text = str(item.get("text") or "")
    paragraph.alignment = _alignment_from_css(item.get("textAlign"))
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)

    line_height = str(item.get("lineHeight") or "").strip().lower()
    if line_height and line_height != "normal":
        paragraph.line_spacing = Pt(_px_to_pt(line_height))

    font = paragraph.font
    font.name = _font_name_from_css(item.get("fontFamily"))
    font.size = Pt(_px_to_pt(item.get("fontSize"), default=14.0))
    font.bold = _weight_from_css(item.get("fontWeight")) >= 600
    font.italic = str(item.get("fontStyle") or "").strip().lower() == "italic"

    parsed_color = _parse_css_color(item.get("color"))
    if parsed_color:
        from pptx.dml.color import RGBColor

        rgb, _alpha = parsed_color
        font.color.rgb = RGBColor(*rgb)


def export_pptx_structured(
    project: DesignerProject,
    pages: Optional[str] = None,
    directory: pathlib.Path | str | os.PathLike | None = None,
) -> bytes:
    """Render browser-positioned text, images, and fills into an editable PPTX."""
    selected = _selected_pages(project, pages)

    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError(
            "python-pptx is required for PPTX export. Run: pip install python-pptx"
        ) from exc
    from pptx.dml.color import RGBColor
    from pptx.util import Emu

    prs = Presentation()
    prs.slide_width = Emu(int(project.canvas_width * 914400 / 96))
    prs.slide_height = Emu(int(project.canvas_height * 914400 / 96))
    blank_layout = prs.slide_layouts[6]

    brand = project.brand
    try:
        bg = RGBColor.from_string(brand.bg_color.lstrip("#")) if brand else RGBColor(0x0F, 0x17, 0x2A)
    except Exception:
        bg = RGBColor(0x0F, 0x17, 0x2A)
    try:
        text_c = RGBColor.from_string(brand.text_color.lstrip("#")) if brand else RGBColor(0xF8, 0xFA, 0xFC)
    except Exception:
        text_c = RGBColor(0xF8, 0xFA, 0xFC)

    for page_index, page in selected:
        slide = prs.slides.add_slide(blank_layout)

        snapshot, screenshots = _collect_rendered_slide_snapshot(project, page.html, page_index=page_index)

        bg_fill = slide.background.fill
        if not _apply_fill(bg_fill, snapshot.get("backgroundColor")):
            bg_fill.solid()
            bg_fill.fore_color.rgb = bg

        items = sorted(snapshot.get("items", []), key=_rendered_item_sort_key)
        for item in items:
            _add_rendered_item_to_slide(slide, item, screenshots)

        if page.notes:
            slide.notes_slide.notes_text_frame.text = page.notes

    buf = io.BytesIO()
    prs.save(buf)
    out_path = describe_export_destination(project, "pptx", pages, mode="structured", directory=directory)
    return _save_bytes(out_path, buf.getvalue(), "Editable PPTX")